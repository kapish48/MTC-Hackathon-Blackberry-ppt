# pip install python-pptx pillow
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image

# ---------- Background ----------
def make_vertical_gradient(width, height, top_rgb, bottom_rgb, path):
    img = Image.new("RGB", (width, height), color=0)
    for y in range(height):
        t = y / (height - 1)
        r = int(top_rgb[0] * (1 - t) + bottom_rgb[0] * t)
        g = int(top_rgb[1] * (1 - t) + bottom_rgb[1] * t)
        b = int(top_rgb[2] * (1 - t) + bottom_rgb[2] * t)
        for x in range(width):
            img.putpixel((x, y), (r, g, b))
    img.save(path)

bg_path = "bb_gradient_bg.png"
logo_path = "logomain.png"  # your white BlackBerry PNG
make_vertical_gradient(1920, 1080, (10, 15, 25), (5, 35, 70), bg_path)

# ---------- Helpers ----------
def apply_background(slide, prs):
    # Background
    slide.shapes.add_picture(bg_path, 0, 0, prs.slide_width, prs.slide_height)

    # Logo bottom-right
    logo_size = Inches(1.2)
    slide.shapes.add_picture(
        logo_path,
        prs.slide_width - logo_size - Inches(0.3),   # right margin
        prs.slide_height - logo_size - Inches(0.3),  # bottom margin
        logo_size,
        logo_size
    )


def style_paragraph(p, font_size=20, bold=False, bullet=True, align=PP_ALIGN.LEFT):
    p.alignment = align
    p.space_after = Pt(12)
    p.space_before = Pt(6)
    p.line_spacing = 1.2

    if bullet:
        p.level = p.level or 0
        if not p.text.strip().startswith("•"):
            p.text = "• " + p.text

    for run in p.runs:
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.name = "Times New Roman" if bold else "Calibri"
        run.font.color.rgb = RGBColor(255, 255, 255)

def add_title_slide(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, prs)

    tb = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12), Inches(2))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = title
    for p in tf.paragraphs:
        style_paragraph(p, font_size=48, bold=True, bullet=False, align=PP_ALIGN.CENTER)

    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(0.7), Inches(3.5), Inches(12), Inches(1.5))
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        tf2.text = subtitle
        for p in tf2.paragraphs:
            style_paragraph(p, font_size=24, bullet=False, align=PP_ALIGN.CENTER)

    return slide

def add_bullets_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, prs)

    tb_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1.5))
    tf_title = tb_title.text_frame
    tf_title.word_wrap = True
    tf_title.text = title
    for p in tf_title.paragraphs:
        style_paragraph(p, font_size=32, bold=True, bullet=False, align=PP_ALIGN.CENTER)

    tb = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.5), Inches(5.5))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.clear()

    def add_paragraph(text, level=0):
        p = tf.add_paragraph() if tf.text != "" else tf.paragraphs[0]
        p.text = text
        p.level = level
        return p

    for item in bullets:
        if isinstance(item, list):
            for sub in item:
                p = add_paragraph(sub, 1)
                style_paragraph(p)
        else:
            p = add_paragraph(item)
            style_paragraph(p)

    return slide

def add_conclusion_slide(prs, title="Conclusion", points=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, prs)

    tb_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12), Inches(1.5))
    tf_title = tb_title.text_frame
    tf_title.word_wrap = True
    tf_title.text = title
    for p in tf_title.paragraphs:
        style_paragraph(p, font_size=40, bold=True, bullet=False, align=PP_ALIGN.CENTER)

    if points:
        tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.5), Inches(4))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.clear()
        for point in points:
            p = tf.add_paragraph()
            p.text = point
            style_paragraph(p, font_size=24, bullet=True, align=PP_ALIGN.LEFT)

    return slide

def add_bullets_slide_custom(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, prs)

    tb_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1.5))
    tf_title = tb_title.text_frame
    tf_title.word_wrap = True
    tf_title.text = title
    for p in tf_title.paragraphs:
        style_paragraph(p, font_size=32, bold=True, bullet=False, align=PP_ALIGN.CENTER)

    tb = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.5), Inches(5.5))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.clear()

    def add_paragraph(text, level=0, bullet=True, font_size=20, bold=False):
        p = tf.add_paragraph() if tf.text != "" else tf.paragraphs[0]
        p.text = text
        p.level = level
        style_paragraph(p, font_size=font_size, bullet=bullet, bold=bold)
        return p

    # Section 1
    add_paragraph("1. The Rise of BYOD (2010–2015):", bullet=False, font_size=22, bold=True)
    add_paragraph("Employees wanted iPhones & Android phones (better apps, touchscreens).", level=1)
    add_paragraph("Companies allowed it due to cost savings and higher employee satisfaction.", level=1)
    add_paragraph("By 2012, 78% of companies supported BYOD.", level=1)

    # Section 2
    add_paragraph("2. Blackberry's mistake - Ignoring BYOD:", bullet=False, font_size=22, bold=True)
    add_paragraph("Assumed enterprises would stick with BlackBerry for security.", level=1)
    add_paragraph("Reality: Apple & Google rapidly improved mobile security and MDM.", level=1)
    add_paragraph("Employees disliked carrying two phones (work BlackBerry + personal iPhone).", level=1)

    add_paragraph("Outcome: rapid migration to iOS/Android, erosion of BlackBerry’s enterprise base.", level=1)

    return slide

# ---------- Build Presentation ----------
prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)

add_title_slide(prs, "BLACKBERRY", "A concise analysis of strategy, product, and market shifts")

add_bullets_slide(prs, "Why did it fail?", [
    "Google’s Android devices quickly followed iPhone with large multitouch screens and no physical keyboard—plus more customization.",
    "BlackBerry believed physical QWERTY keyboards were irreplaceable for professionals.",
    "Focused heavily on productivity, not entertainment or media consumption.",
    "Design philosophy: small screen + keyboard = best professional device."
])

slide3 = add_bullets_slide(prs, "Major three reasons for the failure", [
    "Slow adaptation to touchscreens and modern smartphone trends",
    "Strategic missteps in the app ecosystem",
    "Overconfidence in enterprise market dominance"
])
# Insert chart image below bullet points
chart_path = "Sales.jpg"
slide3.shapes.add_picture(
    chart_path,
    Inches(2),    # x position
    Inches(3.75),  # y position (below bullet points)
    Inches(5.5),    # width
    Inches(3.5)   # height
)

add_bullets_slide(prs, "1) Slow Adaptation to Touchscreens & Modern Smartphone Trends", [
    "Relied on physical keyboards while iPhone & Android embraced large touchscreens.",
    "Underestimated consumer demand for multimedia, apps, and gesture-based navigation.",
    "First touchscreen phone (BlackBerry Storm, 2008) had poor performance and clunky “SurePress” screen.",
    "Delayed response allowed rivals to capture both consumer and enterprise markets.",
    "Developers shifted focus to iOS/Android, widening BlackBerry’s innovation gap.",
    "Failure to redesign OS for touch led to an awkward, outdated user experience."
])

slide5 = add_bullets_slide(prs, "BlackBerry Storm (“SurePress”)", [
    "First touchscreen phone (BlackBerry Storm, 2008): tried to mimic pressing a real key.",
    "Major Flaws: ",
    [
        "Touch response lag",
        "Clicks felt unnatural and tiring",
        "OS not designed for touch → awkward navigation",
        "App support was weak vs Apple App Store"
    ],
    "Result: negative reviews, high return rates, and product failure."
])
# Insert chart image below bullet points
chart_path = "storm.png"
slide5.shapes.add_picture(
    chart_path,
    Inches(6.5),    # x position
    Inches(1.5),  # y position (below bullet points)
    Inches(7.75),    # width
    Inches(5.75)   # height
)
add_bullets_slide(prs, "2) BlackBerry’s Strategic Missteps in the App Ecosystem", [
    "Prioritized secure email and BBM, neglecting a modern, user-friendly app store.",
    "Did not recognize early enough that a wide variety of mobile apps was becoming crucial.",
    "Provided limited tools/incentives for third-party developers → slow innovation & fewer apps.",
    "Overlooked the speed of iOS/Android platform evolution to attract users and developers."
])

add_bullets_slide(prs, "Consequences of Underestimating the App Ecosystem", [
    "Most developers shifted to iOS and Android → a “BlackBerry app gap.”",
    "Users faced restricted choices for social, entertainment, productivity, and lifestyle apps.",
    "Frustrated by missing/outdated apps, many migrated to competitors with richer ecosystems.",
    "This app shortfall drove loss of relevance and a steep decline in global market share."
])

add_bullets_slide(prs, "3) Overconfidence in Enterprise Market Dominance", [
    "2009 Peak: Controlled 50% of the U.S. smartphone market, with 20% globally (Forbes).",
    "Government/Corporate Reliance: Used by 90% of Fortune 500 companies and governments (BBC).",
    "Assumption: Believed businesses would prioritize security over employee preferences (WSJ).",
    "Why it failed:",
    [
        "Ignored BYOD: By 2012, 78% of companies allowed employees to use personal devices (Gartner).",
        "Slow to Adapt: BlackBerry’s CEO called BYOD a 'passing trend' in 2012 (The Verge).",
        "Security Myth: iOS/Android adopted MDM (Mobile Device Management), matching BlackBerry’s security (CNBC)."
    ]
])

add_bullets_slide_custom(prs, "How BYOD Killed BlackBerry")

add_conclusion_slide(prs, "Final Takeaways", [
    "BlackBerry’s fall was not due to one mistake, but a series of missed shifts.",
    "Failure to adapt to consumer trends and app ecosystems proved fatal.",
    "Lesson: In technology, adaptability and user experience matter as much as security and performance."
])

# Save
prs.save("blackberry_with_logo.pptx")
print("Saved blackberry_with_logo.pptx")
